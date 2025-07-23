import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
import pytesseract
from PIL import Image
from io import BytesIO
from langchain.text_splitter import RecursiveCharacterTextSplitter

class IngestionAgent:
    def parse_and_preprocess(self, files):
        raw_texts = []

        for file in files:
            if file.type == "application/pdf":
                raw_texts.extend(self._parse_pdf(file))
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                raw_texts.extend(self._parse_docx(file))
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                raw_texts.extend(self._parse_pptx(file))
            elif file.type == "text/plain":
                raw_texts.append(file.read().decode("utf-8"))
            elif file.type == "text/csv":
                df = pd.read_csv(file)
                raw_texts.append(df.to_string())

        all_chunks = []
        for text in raw_texts:
            all_chunks.extend(self._split_text(text))

        return all_chunks

    def _split_text(self, text, chunk_size=500, chunk_overlap=50):
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        return splitter.split_text(text)

    def _parse_pdf(self, file):
        file_bytes = file.read()
        file.seek(0)
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        text_chunks = []

        for page in doc:
            text_chunks.append(page.get_text())

            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                image = Image.open(BytesIO(image_bytes))
                text_chunks.append("[IMAGE OCR TEXT] " + pytesseract.image_to_string(image))

        return text_chunks

    def _parse_docx(self, file):
        doc = docx.Document(file)
        text_chunks = [para.text for para in doc.paragraphs if para.text.strip()]

        for rel in doc.part._rels:
            if "image" in doc.part._rels[rel].target_ref:
                image_part = doc.part._rels[rel].target_part
                img = Image.open(BytesIO(image_part.blob))
                text_chunks.append("[IMAGE OCR TEXT] " + pytesseract.image_to_string(img))

        return text_chunks

    def _parse_pptx(self, file):
        prs = pptx.Presentation(file)
        text_chunks = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_chunks.append(shape.text)
                elif shape.shape_type == 13:  # Picture
                    image = shape.image
                    img = Image.open(BytesIO(image.blob))
                    text_chunks.append("[IMAGE OCR TEXT] " + pytesseract.image_to_string(img))

        return text_chunks
